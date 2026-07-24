"""Local text extraction from attachment files.

Extracts text from PDF, Word, PowerPoint, Excel, images (OCR),
.eml, .rpmsg, and plain text files. No API calls — all local.
"""

import html
import os
import re
from email import policy
from email.message import EmailMessage
from email.parser import BytesParser
from pathlib import Path

# Maximum characters to store per file
MAX_TEXT_CHARS = 100_000
# Minimum characters to consider a successful extraction
MIN_TEXT_CHARS = 50

# MIME types we skip entirely (video, audio, archives, Outlook artifacts)
SKIP_MIME_TYPES = {
    "video/mp4",
    "audio/mpeg",
    "audio/x-wav",
    "audio/wav",
    "application/zip",
    "application/x-rar-compressed",
    "application/x-7z-compressed",
    "application/gzip",
}

SKIP_EXTENSIONS = {".mp4", ".mp3", ".wav", ".zip", ".rar", ".7z", ".gz", ".mso", ".wmz"}


def extract_text_from_file(file_path: str, mime_type: str) -> dict:
    """Extract text from a file based on its MIME type.

    Returns dict with keys: text, method, status, error.
    status is one of: 'extracted', 'partial', 'failed', 'skipped'.
    """
    ext = Path(file_path).suffix.lower()

    # Skip unsupported types
    if mime_type in SKIP_MIME_TYPES or ext in SKIP_EXTENSIONS:
        return {"text": None, "method": None, "status": "skipped", "error": None}

    # Check file exists
    if not os.path.isfile(file_path):
        return {
            "text": None,
            "method": None,
            "status": "failed",
            "error": f"File not found: {file_path}",
        }

    try:
        if mime_type == "application/pdf" or ext == ".pdf":
            return _extract_pdf(file_path)
        elif (
            mime_type
            in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",)
            or ext == ".docx"
        ):
            return _extract_docx(file_path)
        elif mime_type == "application/msword" or ext == ".doc":
            return _extract_doc(file_path)
        elif (
            mime_type
            in ("application/vnd.openxmlformats-officedocument.presentationml.presentation",)
            or ext == ".pptx"
        ):
            return _extract_pptx(file_path)
        elif ext == ".xlsb":
            return _extract_xlsb(file_path)
        elif ext == ".xls":
            return _extract_xls(file_path)
        elif (
            mime_type
            in (
                "application/vnd.ms-excel",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
            or ext == ".xlsx"
        ):
            return _extract_excel(file_path)
        elif (
            mime_type
            and mime_type.startswith("image/")
            or ext in (".png", ".jpg", ".jpeg", ".gif", ".tiff", ".tif", ".bmp", ".jfif")
        ):
            return _extract_image_ocr(file_path)
        elif mime_type == "message/rfc822" or ext == ".eml":
            return _extract_eml(file_path)
        elif mime_type == "application/encrypted" or ext == ".rpmsg":
            return _extract_rpmsg(file_path)
        elif mime_type in ("text/plain", "text/csv", "text/markdown") or ext in (
            ".txt",
            ".csv",
            ".md",
        ):
            return _extract_plain_text(file_path)
        elif mime_type == "text/html" or ext in (".html", ".htm"):
            return _extract_html(file_path)
        else:
            return {
                "text": None,
                "method": None,
                "status": "skipped",
                "error": f"Unsupported type: {mime_type} ({ext})",
            }
    except Exception as e:
        return {
            "text": None,
            "method": None,
            "status": "failed",
            "error": f"{type(e).__name__}: {str(e)[:500]}",
        }


def _extract_pdf(path: str) -> dict:
    """Extract text from PDF using PyMuPDF, with OCR fallback."""
    import fitz

    doc = fitz.open(path)
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()

    text = "\n\n".join(pages)
    text = _truncate(text)

    # If very little text extracted, fall back to per-page OCR (scanned PDF).
    # _extract_image_ocr cannot read PDFs (Image.open fails); _ocr_pdf_pages
    # renders each page via fitz.get_pixmap before OCR-ing. See B3 spec.
    if len(text.strip()) < MIN_TEXT_CHARS:
        ocr_result = _ocr_pdf_pages(path)
        if ocr_result["status"] == "extracted" and len(ocr_result["text"] or "") > len(
            text.strip()
        ):
            return ocr_result
        if _apply_noise_filter(text):
            return {
                "text": None,
                "method": "pymupdf",
                "status": "skipped",
                "error": "Insufficient text extracted",
            }

    return {"text": text, "method": "pymupdf", "status": "extracted", "error": None}


def _extract_docx(path: str) -> dict:
    """Extract text from .docx using python-docx."""
    from docx import Document

    doc = Document(path)
    parts = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # Extract table content
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    text = _truncate("\n".join(parts))
    if _apply_noise_filter(text):
        return {
            "text": None,
            "method": "python-docx",
            "status": "skipped",
            "error": "Insufficient text extracted",
        }
    return {"text": text, "method": "python-docx", "status": "extracted", "error": None}


def _extract_doc(path: str) -> dict:
    """Extract text from legacy .doc files using textutil (macOS)."""
    import subprocess

    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", path],
            capture_output=True,
            text=True,
            timeout=30,
        )
        text = _truncate(result.stdout)
        if _apply_noise_filter(text):
            return {
                "text": None,
                "method": "textutil",
                "status": "skipped",
                "error": "Insufficient text extracted",
            }
        return {
            "text": text,
            "method": "textutil",
            "status": "extracted",
            "error": None,
        }
    except (subprocess.TimeoutExpired, FileNotFoundError) as e:
        return {"text": None, "method": "textutil", "status": "failed", "error": str(e)}


def _extract_pptx(path: str) -> dict:
    """Extract text from PowerPoint using python-pptx."""
    from pptx import Presentation

    prs = Presentation(path)
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_text.append(" | ".join(cells))
        if slide_text:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes] {notes}")

    text = _truncate("\n\n".join(parts))
    if _apply_noise_filter(text):
        return {
            "text": None,
            "method": "python-pptx",
            "status": "skipped",
            "error": "Insufficient text extracted",
        }
    return {"text": text, "method": "python-pptx", "status": "extracted", "error": None}


def _extract_excel(path: str) -> dict:
    """Extract headers + first 50 rows per sheet from .xlsx files via openpyxl.

    .xlsb and .xls are dispatched to dedicated parsers (_extract_xlsb,
    _extract_xls) before reaching this function — see extract_text_from_file.
    """
    import openpyxl

    parts = []
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)

    max_sheets = 20
    for _sheet_idx, sheet_name in enumerate(wb.sheetnames[:max_sheets]):
        ws = wb[sheet_name]
        rows_text = []
        row_count = 0
        for row in ws.iter_rows(max_row=51, values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows_text.append(" | ".join(c for c in cells if c.strip()))
            row_count += 1
            if row_count >= 51:
                break

        if rows_text:
            parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows_text))

    wb.close()

    text = _truncate("\n\n".join(parts))
    if _apply_noise_filter(text):
        return {
            "text": None,
            "method": "openpyxl",
            "status": "skipped",
            "error": "Insufficient text extracted",
        }
    return {"text": text, "method": "openpyxl", "status": "extracted", "error": None}


def _extract_xls(path: str) -> dict:
    """Extract headers + first 50 rows per sheet from legacy .xls files.

    xlrd 2.0+ dropped .xlsx support and now handles only the legacy BIFF
    .xls format — perfect fit for our case. Mirrors _extract_excel's shape
    (sheet headers, max 20 sheets, max 51 rows per sheet) so downstream
    LLM extraction sees the same structure regardless of source format.
    """
    import xlrd

    try:
        wb = xlrd.open_workbook(path)
    except Exception as e:
        return {
            "text": None,
            "method": "xlrd",
            "status": "failed",
            "error": f"xlrd open failed: {e}",
        }

    parts = []
    max_sheets = 20
    for sheet in list(wb.sheets())[:max_sheets]:
        rows_text = []
        for row_idx in range(min(sheet.nrows, 51)):
            cells = [str(c) if c is not None else "" for c in sheet.row_values(row_idx)]
            if any(c.strip() for c in cells):
                rows_text.append(" | ".join(c for c in cells if c.strip()))
        if rows_text:
            parts.append(f"--- Sheet: {sheet.name} ---\n" + "\n".join(rows_text))

    text = _truncate("\n\n".join(parts))
    if _apply_noise_filter(text):
        return {
            "text": None,
            "method": "xlrd",
            "status": "skipped",
            "error": "Insufficient text extracted",
        }
    return {"text": text, "method": "xlrd", "status": "extracted", "error": None}


def _extract_xlsb(path: str) -> dict:
    """Extract headers + first 50 rows per sheet from .xlsb (Excel binary) files.

    pyxlsb is the only mature Python reader for the .xlsb format. It loads
    the whole workbook (no read_only mode like openpyxl), but our per-sheet
    + per-row caps below bound the per-file work.
    """
    import pyxlsb

    parts = []
    max_sheets = 20

    try:
        with pyxlsb.open_workbook(path) as wb:
            for sheet_name in list(wb.sheets)[:max_sheets]:
                rows_text = []
                with wb.get_sheet(sheet_name) as sheet:
                    for row_idx, row in enumerate(sheet.rows()):
                        if row_idx >= 51:
                            break
                        cells = [
                            str(cell.v) if cell is not None and cell.v is not None else ""
                            for cell in row
                        ]
                        if any(c.strip() for c in cells):
                            rows_text.append(" | ".join(c for c in cells if c.strip()))
                if rows_text:
                    parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows_text))
    except Exception as e:
        # Many real-world files have a .xlsb extension but are actually OLE
        # compound documents (legacy .xls magic d0cf11e0...). pyxlsb errors
        # with "File is not a zip file" on those. Detect and fall through to
        # xlrd, which handles the OLE BIFF format. If xlrd also fails, return
        # the original pyxlsb error since that's the user-facing classification.
        if "not a zip file" in str(e).lower():
            xls_result = _extract_xls(path)
            if xls_result["status"] == "extracted":
                # Tag method so downstream debugging can see the extension/format mismatch.
                xls_result["method"] = "xlrd (fallback from .xlsb)"
                return xls_result
            # OLE compound document with no Excel workbook stream — likely a
            # custom export format (ACME financial tools, etc.) that wraps
            # binary content in OLE but has no spreadsheet payload. No library
            # will rescue this as Excel data; classify as skipped (not failed)
            # so future retries don't waste cycles on it.
            if "find workbook" in (xls_result.get("error") or "").lower():
                return {
                    "text": None,
                    "method": "pyxlsb→xlrd",
                    "status": "skipped",
                    "error": "OLE compound document, no Excel workbook stream — likely custom export format, not standard Excel",
                }
        return {
            "text": None,
            "method": "pyxlsb",
            "status": "failed",
            "error": f"pyxlsb open failed: {e}",
        }

    text = _truncate("\n\n".join(parts))
    if _apply_noise_filter(text):
        return {
            "text": None,
            "method": "pyxlsb",
            "status": "skipped",
            "error": "Insufficient text extracted",
        }
    return {"text": text, "method": "pyxlsb", "status": "extracted", "error": None}


def _extract_image_ocr(path: str) -> dict:
    """Extract text from images using Tesseract OCR."""
    try:
        import io

        import pytesseract
        from PIL import Image
    except ImportError:
        return {
            "text": None,
            "method": "ocr",
            "status": "failed",
            "error": "pytesseract or Pillow not installed",
        }

    # pytesseract rejects any image whose PIL .format isn't in its allowlist
    # (JPEG, PNG, GIF, BMP, TIFF, WEBP, PPM). Phone-camera JPEGs are often
    # encoded as MPO (Multi-Picture Object) — a JPEG container with multiple
    # frames — and tesseract raises TypeError: Unsupported image format/type.
    # Re-encode through PNG to drop the container and yield a clean format.
    _TESSERACT_SAFE_FORMATS = {"JPEG", "PNG", "GIF", "BMP", "TIFF", "WEBP", "PPM"}

    try:
        img = Image.open(path)
        if img.format not in _TESSERACT_SAFE_FORMATS:
            buf = io.BytesIO()
            img.convert("RGB").save(buf, format="PNG")
            buf.seek(0)
            img = Image.open(buf)
        text = pytesseract.image_to_string(img, lang="eng+ell")
        text = _truncate(text)

        if _apply_noise_filter(text):
            return {
                "text": None,
                "method": "ocr",
                "status": "skipped",
                "error": "OCR returned insufficient text",
            }
        return {"text": text, "method": "ocr", "status": "extracted", "error": None}
    except Exception as e:
        return {
            "text": None,
            "method": "ocr",
            "status": "failed",
            "error": f"{type(e).__name__}: {str(e)[:200]}",
        }


def _ocr_pdf_pages(path: str, max_pages: int = 30) -> dict:
    """Render each PDF page to a PIL image and OCR it. Used as fallback for
    scanned PDFs when PyMuPDF text extraction yields too little content.

    200 DPI is the sweet spot for printed text — enough resolution for
    Tesseract to recognize Greek + English glyphs reliably, low enough that
    a typical 2-page scan completes in 3-10 seconds. Hard-cap at 30 pages
    to bound worst-case time on big documents (most ACME scans are <10).
    """
    import io

    import fitz
    import pytesseract
    from PIL import Image

    try:
        doc = fitz.open(path)
        pages_text = []
        for page_num, page in enumerate(doc):
            if page_num >= max_pages:
                break
            pix = page.get_pixmap(dpi=200)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            page_text = pytesseract.image_to_string(img, lang="eng+ell")
            if page_text.strip():
                pages_text.append(page_text.strip())
        doc.close()
    except Exception as e:
        return {
            "text": None,
            "method": "pymupdf+tesseract",
            "status": "failed",
            "error": f"{type(e).__name__}: {str(e)[:200]}",
        }

    text = _truncate("\n\n".join(pages_text))
    if _apply_noise_filter(text):
        return {
            "text": None,
            "method": "pymupdf+tesseract",
            "status": "skipped",
            "error": "OCR returned insufficient text",
        }
    return {
        "text": text,
        "method": "pymupdf+tesseract",
        "status": "extracted",
        "error": None,
    }


def _extract_eml(path: str) -> dict:
    """Parse .eml files to extract headers and body text."""
    with open(path, "rb") as f:
        msg: EmailMessage = BytesParser(policy=policy.default).parse(f)  # type: ignore[assignment]

    parts = []
    parts.append(f"Subject: {msg.get('subject', 'N/A')}")
    parts.append(f"From: {msg.get('from', 'N/A')}")
    parts.append(f"Date: {msg.get('date', 'N/A')}")
    parts.append("")

    body = msg.get_body(preferencelist=("plain", "html"))
    if body:
        content = body.get_content()
        if body.get_content_type() == "text/html":
            content = _strip_html_tags(content)
        parts.append(content)

    text = _truncate("\n".join(parts))
    if _apply_noise_filter(text):
        return {
            "text": None,
            "method": "email_parser",
            "status": "skipped",
            "error": "Insufficient text extracted",
        }
    return {
        "text": text,
        "method": "email_parser",
        "status": "extracted",
        "error": None,
    }


def _extract_rpmsg(path: str) -> dict:
    """Best-effort metadata extraction from encrypted .rpmsg files."""
    try:
        import compoundfiles
    except ImportError:
        return {
            "text": None,
            "method": "compoundfiles",
            "status": "failed",
            "error": "compoundfiles not installed",
        }

    try:
        doc = compoundfiles.CompoundFileReader(path)
        parts = []

        for entry in doc.root:
            name = entry.name
            if entry.is_file:
                try:
                    data = doc.open(entry).read()
                    for encoding in ("utf-8", "utf-16-le", "latin-1"):
                        try:
                            decoded = data.decode(encoding)
                            printable = "".join(
                                c for c in decoded if c.isprintable() or c in "\n\r\t"
                            )
                            if len(printable) > 10:
                                parts.append(f"[{name}] {printable[:2000]}")
                                break
                        except (UnicodeDecodeError, ValueError):
                            continue
                except Exception:
                    continue

        doc.close()

        text = _truncate("\n".join(parts))
        if not text.strip():
            return {
                "text": None,
                "method": "compoundfiles",
                "status": "partial",
                "error": "Encrypted content — metadata only",
            }
        return {
            "text": text,
            "method": "compoundfiles",
            "status": "partial",
            "error": None,
        }
    except Exception as e:
        return {
            "text": None,
            "method": "compoundfiles",
            "status": "failed",
            "error": f"{type(e).__name__}: {str(e)[:200]}",
        }


def _extract_plain_text(path: str) -> dict:
    """Read plain text, CSV, or markdown files directly."""
    for encoding in ("utf-8", "latin-1", "cp1253"):
        try:
            with open(path, encoding=encoding) as f:
                text = f.read()
            text = _truncate(text)
            if _apply_noise_filter(text):
                return {
                    "text": None,
                    "method": "direct_read",
                    "status": "skipped",
                    "error": "Insufficient text",
                }
            return {
                "text": text,
                "method": "direct_read",
                "status": "extracted",
                "error": None,
            }
        except UnicodeDecodeError:
            continue
    return {
        "text": None,
        "method": "direct_read",
        "status": "failed",
        "error": "Could not decode file with any supported encoding",
    }


def _extract_html(path: str) -> dict:
    """Read HTML files and strip tags."""
    for encoding in ("utf-8", "latin-1", "cp1253"):
        try:
            with open(path, encoding=encoding) as f:
                raw = f.read()
            text = _strip_html_tags(raw)
            text = _truncate(text)
            if _apply_noise_filter(text):
                return {
                    "text": None,
                    "method": "direct_read",
                    "status": "skipped",
                    "error": "Insufficient text",
                }
            return {
                "text": text,
                "method": "direct_read",
                "status": "extracted",
                "error": None,
            }
        except UnicodeDecodeError:
            continue
    return {
        "text": None,
        "method": "direct_read",
        "status": "failed",
        "error": "Could not decode file",
    }


def _strip_html_tags(raw_html: str) -> str:
    """Remove HTML tags and decode entities."""
    text = re.sub(r"<style[^>]*>.*?</style>", "", raw_html, flags=re.DOTALL)
    text = re.sub(r"<script[^>]*>.*?</script>", "", text, flags=re.DOTALL)
    text = re.sub(r"<[^>]+>", " ", text)
    text = html.unescape(text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _truncate(text: str, max_chars: int = MAX_TEXT_CHARS) -> str:
    """Truncate text to max_chars."""
    if len(text) > max_chars:
        return text[:max_chars]
    return text


def _apply_noise_filter(text: str) -> bool:
    """Return True if text should be filtered out (noise).

    Filters: too short (<50 chars) or >90% non-alphanumeric.
    """
    stripped = (text or "").strip()
    if len(stripped) < MIN_TEXT_CHARS:
        return True
    alnum = sum(1 for c in stripped if c.isalnum())
    if alnum / len(stripped) < 0.1:
        return True
    return False
